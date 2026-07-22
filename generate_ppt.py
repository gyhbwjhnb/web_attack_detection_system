"""
生成项目汇报 PPT —— 基于 python-pptx。
输出: docs/项目汇报.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUTPUT = "docs/项目汇报.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ==================== 辅助函数 ====================

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xBF, 0xA5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)
CARD_BG = RGBColor(0x25, 0x25, 0x40)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_card(slide, left, top, width, height, title, content_lines, icon="●"):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)

    p = tf.paragraphs[0]
    p.text = f"{icon} {title}"
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT
    p.font.bold = True
    p.font.name = "Microsoft YaHei"

    for line in content_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(4)


# ==================== Slide 1: 封面 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

add_text_box(slide, 1.5, 1.5, 10, 1.2,
             "网络攻击检测系统 (NIDS)",
             font_size=48, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 2, 2.8, 9, 0.8,
             "基于特征匹配 + 异常行为分析的双引擎入侵检测系统",
             font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 4, 3.8, 5, 0.6,
             "版本 v2.0  |  2026年7月",
             font_size=16, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

# 分隔线
shape = slide.shapes.add_shape(
    1, Inches(4.5), Inches(4.7), Inches(4.3), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# 底部特性
add_text_box(slide, 1.5, 5.0, 10, 1.5,
             "Python 3  |  tkinter GUI  |  10 个检测器插件  |  双引擎并行  |  规则热加载  |  智能评分",
             font_size=14, color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER)


# ==================== Slide 2: 目录 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 5, 0.7, "目录", font_size=36, color=ACCENT, bold=True)

toc_items = [
    ("01", "项目背景与目标"),
    ("02", "系统架构设计"),
    ("03", "模块详细设计"),
    ("04", "检测能力总览"),
    ("05", "智能增强特性 v2.0"),
    ("06", "可扩展性设计"),
    ("07", "GUI 功能展示"),
    ("08", "测试与部署"),
    ("09", "与 Snort3 对比分析"),
    ("10", "后续规划"),
]

for i, (num, title) in enumerate(toc_items):
    y = 1.5 + i * 0.55
    add_text_box(slide, 1.5, y, 1, 0.5, num, font_size=28, color=ACCENT, bold=True)
    add_text_box(slide, 2.8, y + 0.05, 8, 0.5, title, font_size=22, color=WHITE)


# ==================== Slide 3: 项目背景 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.3, 5, 0.7, "项目背景与目标", font_size=36, color=ACCENT, bold=True)

add_card(slide, 0.8, 1.3, 5.8, 5.5, "项目背景", [
    "网络攻击日益复杂化、自动化，传统防火墙无法应对",
    "Snort/Suricata 功能强大但配置复杂，学习曲线陡峭",
    "现有 NIDS 缺乏行为基线学习和智能告警增强能力",
    "教学场景需要可演示、可定制、可扩展的检测平台",
    "需要 Python 原生方案，降低部署和维护成本",
], icon="🎯")

add_card(slide, 7.2, 1.3, 5.8, 5.5, "设计目标", [
    "双引擎检测：特征匹配 + 异常行为分析并行工作",
    "插件化架构：新增检测器仅需 3 步，不动核心代码",
    "零依赖 GUI：使用 Python 内置 tkinter 无需额外安装",
    "规则热加载：拖放 JSON 文件到目录即生效，无需重启",
    "智能增强：置信度评分 + 时段检测 + 自学习误报抑制",
    "离线可用：支持 PCAP 回放和模拟数据，无需 Npcap",
], icon="🚀")


# ==================== Slide 4: 系统架构 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.3, 5, 0.7, "系统架构设计", font_size=36, color=ACCENT, bold=True)

# 架构图用文本框模拟
add_text_box(slide, 3.5, 1.2, 6, 0.5,
             "┌─────────── MessageBus(事件总线) ───────────┐",
             font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_card(slide, 0.8, 2.0, 2.8, 4.5, "模块1: Capture", [
    "三种数据源:",
    "  在线网卡抓包",
    "  离线 PCAP 回放",
    "  内存模拟注入",
    "协议解析:",
    "  HTTP/DNS/TLS/ICMP",
    "  ARP/SMB 逐层解析",
    "输出统一 TrafficRecord",
], icon="📡")

add_card(slide, 4.0, 2.0, 2.8, 4.5, "模块2: 特征匹配", [
    "Aho-Corasick 多模式",
    "  匹配（单次扫描）",
    "URL 解码对抗绕过",
    "三级过滤:",
    "  协议 → 端口 → 正则",
    "暴力破解频次统计",
    "告警去重（窗口期）",
    "规则热加载监听",
], icon="🔍")

add_card(slide, 7.2, 2.0, 2.8, 4.5, "模块3: 异常检测", [
    "EMA 自适应行为基线",
    "10 个检测器插件:",
    "  端口扫描/暴力破解",
    "  DDoS/异常外联/横向扩散",
    "  带宽异常/行为突变",
    "  时段检测/置信度评分",
    "  自学习误报抑制",
    "攻击链构建",
], icon="🧠")

add_card(slide, 10.4, 2.0, 2.5, 4.5, "模块4: GUI", [
    "实时告警展示",
    "流量监控（按协议着色）",
    "三维组合筛选",
    "统计面板",
    "HTML/CSV/TXT 导出",
    "白名单管理",
    "6 标签页帮助系统",
], icon="🖥️")


# ==================== Slide 5: 数据流 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.3, 5, 0.7, "数据流与通信机制", font_size=36, color=ACCENT, bold=True)

add_card(slide, 0.8, 1.3, 5.8, 2.2, "MessageBus 6 个标准事件", [
    "traffic_record  →  模块1发布，模块2/3/4订阅",
    "signature_alert →  模块2发布特征匹配告警",
    "anomaly_alert   →  模块3发布异常检测告警",
    "attack_chain    →  模块3发布 Kill Chain 事件",
    "statistics      →  模块3发布统计摘要",
    "config_change   →  模块4发布配置变更通知",
], icon="📨")

add_card(slide, 7.2, 1.3, 5.8, 2.2, "处理流水线", [
    "1. 模块1: 抓包 → parse_packet() → TrafficRecord",
    "2. publish(traffic_record) → 模块2 + 模块3 并行",
    "3. 模块2: 三级过滤 → 模块1 中的 Aho-Corasick → Alert",
    "4. 模块3: 10 个检测器遍历 → 后处理 → Alert",
    "5. publish(signature_alert / anomaly_alert) → GUI",
], icon="⚙️")

add_card(slide, 0.8, 4.0, 12.1, 3.0, "关键设计决策", [
    "模块间完全解耦：无 import 依赖，仅通过 MessageBus 通信",
    "并行处理：模块2(特征匹配)和模块3(异常检测)同时运行，无先后依赖",
    "优先级排序：检测器按 priority 字段升序执行，后处理阶段 (80-90) 在最末",
    "线程安全：所有检测器共享 self._lock，通过 set_context() 注入",
    "路径无关性：相对路径自动解析为项目根目录（signature_engine._project_root）",
], icon="💡")


# ==================== Slide 6: 检测能力 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.3, 5, 0.7, "检测能力总览", font_size=36, color=ACCENT, bold=True)

add_card(slide, 0.8, 1.3, 5.8, 5.5, "特征匹配检测 (20+ 类攻击)", [
    "SQL注入: UNION SELECT / 报错注入 / 盲注 / 堆叠查询",
    "XSS: <script> 反射 / img onerror / DOM 注入",
    "命令注入: 管道符 / 反引号 / $(...) 执行",
    "文件包含: LFI (/etc/passwd) / RFI (远程URL)",
    "WebShell: ASP 执行 / PHP eval / JSP Runtime",
    "后门: 一句话木马 / Cobalt Strike Beacon",
    "协议: DNS隧道 / Shadowsocks / 反弹Shell",
    "信息泄露: .env / phpinfo / .git 暴露",
    "CVE: CVE-2021-44228 Log4j 通用模式",
], icon="🛡️")

add_card(slide, 7.2, 1.3, 5.8, 5.5, "异常行为检测 (10 个插件)", [
    "端口扫描:  单IP → 多端口  → 超过阈值 → 告警",
    "暴力破解:  短时间  → 密集连接敏感端口 → 加分",
    "异常外联:  内网 → 陌生外网IP + 非标准端口",
    "SYN Flood:  SYN包密度偏离基线 → DDoS 判定",
    "带宽异常:  流量突发 > 基线 × 3.0 → 告警",
    "横向扩散:  1个IP → 多个内网主机 + 跨子网",
    "行为突变:  4维度(连接数/对端/端口/SYN比)",
    "时段检测:  24h EMA 基线，凌晨异常自动提权",
], icon="📊")


# ==================== Slide 7: 智能增强 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.3, 5, 0.7, "智能增强特性 v2.0", font_size=36, color=ACCENT, bold=True)

add_card(slide, 0.8, 1.3, 3.9, 5.5, "S6: 置信度评分", [
    "动态计算 0.0 - 1.0 置信分",
    "7 种攻击类型各有公式:",
    "port_scan: 端口数×0.01",
    "brute_force: 密度 + 端口加分",
    "syn_flood: SYN/基线 ×0.05",
    "outbound: 端口恶意加分",
    "lateral: 主机数×0.03",
    "bandwidth: 偏离倍数×0.2",
    "behavior: 最大偏离×0.2",
    "后处理阶段 priority=85",
], icon="📈")

add_card(slide, 5.1, 1.3, 3.9, 5.5, "S7: 时段检测", [
    "每主机维护 24h EMA 基线",
    "均值 + 2σ 触发告警",
    "仅凌晨 0:00-6:00 生效",
    "置信度自动提升:",
    "  2σ → +0.15",
    "  5σ → +0.25",
    "告警标记:",
    '  "[夜间异常时段(03:00)]"',
    "学习数小时自动稳定",
], icon="🕐")

add_card(slide, 9.4, 1.3, 3.9, 5.5, "S8: 误报抑制", [
    "用户右键标记\"误报\"",
    "提取签名:",
    "  src_ip|attack_type|dst_port",
    "持久化到 JSON 文件",
    "累计 2 次标记 → 自动",
    "  后续同模式 -0.35 置信",
    "重启后自动加载",
    "越用越精准",
], icon="🎓")


# ==================== Slide 8: 可扩展性 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.3, 8, 0.7, "可扩展性设计 — IDetector 插件接口", font_size=36, color=ACCENT, bold=True)

add_card(slide, 0.8, 1.3, 7.0, 5.5, "IDetector 抽象基类", [
    "name: str          — 唯一标识（下划线命名）",
    "category: str      — 分类标签（scan/brute/exfil 等）",
    "priority: int      — 执行顺序（10-50主检测 / 80-90后处理）",
    "enabled: bool      — 开关（默认 True）",
    "",
    "必须实现:",
    "  process(record, now) → List[Alert]   核心检测方法",
    "",
    "可选实现:",
    "  post_process_batch(alerts) → List[Alert]  后处理阶段",
    "",
    "引擎自动注入:",
    "  self._hosts / self._baselines / self._config / self._lock",
    "  self._count_recent() / self._make_alert()",
], icon="🔌")

add_card(slide, 8.2, 1.3, 4.7, 2.5, "新增检测器 3 步", [
    "Step 1: detectors/ 下新建 .py 文件",
    "Step 2: class MyDetector(IDetector):",
    "           name = \"my_detector\"",
    "           def process(self, r, t): ...",
    "Step 3: __init__.py ← 加一行 MyDetector()",
    "",
    "✓ 不改 anomaly_engine.py",
    "✓ 自动注入上下文，按 priority 排序",
], icon="📝")

add_card(slide, 8.2, 4.2, 4.7, 2.5, "规则热加载", [
    "data/rules/ 目录下任意 .json 文件",
    "后台线程 3 秒轮询",
    "文件变化 → 自动 reload",
    "控制台输出确认信息",
    "",
    "GUI → 导入规则按钮 → 一键导入",
    "帮助 → 自定义规则 → 格式参考",
], icon="🔄")


# ==================== Slide 9: GUI 功能 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.3, 5, 0.7, "GUI 功能展示", font_size=36, color=ACCENT, bold=True)

features_left = [
    ("告警列表", "实时展示，按严重度5级着色（蓝→红）"),
    ("流量列表", "正常流量监控，按协议4色标记"),
    ("筛选栏", "IP/协议/端口三维组合筛选，下拉建议"),
    ("统计面板", "总数/攻击类型分布/严重度分布/攻击链"),
    ("置信度", "每条告警动态 0.00-1.00 评分"),
    ("右键菜单", "标记状态（待处理/已确认/误报）+ 备注"),
]
features_right = [
    ("白名单", "GUI 即时添加，引擎自动静默"),
    ("导出报告", "HTML(含分布图)/CSV/TXT 三种格式"),
    ("流量下载", "双击流量行 → 原始载荷 .bin/.json"),
    ("导入规则", "工具栏按钮 + 菜单双入口"),
    ("帮助系统", "6 标签页（快速开始/界面/筛选/白名单/导出/规则）"),
    ("状态栏", "运行时间 + 最近事件提示"),
]

y = 1.3
for title, desc in features_left:
    add_text_box(slide, 1.0, y, 2.5, 0.4, title, font_size=16, color=ACCENT, bold=True)
    add_text_box(slide, 3.5, y, 4, 0.4, desc, font_size=14, color=LIGHT_GRAY)
    y += 0.9

y = 1.3
for title, desc in features_right:
    add_text_box(slide, 8.0, y, 2.5, 0.4, title, font_size=16, color=ACCENT, bold=True)
    add_text_box(slide, 10.5, y, 3, 0.4, desc, font_size=14, color=LIGHT_GRAY)
    y += 0.9


# ==================== Slide 10: 与 Snort3 对比 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.3, 5, 0.7, "与 Snort3 对比分析", font_size=36, color=ACCENT, bold=True)

add_card(slide, 0.8, 1.3, 5.8, 5.5, "Snort3 的优势", [
    "语言: C++，性能极高",
    "数据面: DPDK 零拷贝包处理",
    "规则语法: 自有 DSL，生态成熟",
    "规则库: 数万条社区维护规则",
    "集成: 可嵌入 Suricata/Barnyard2",
    "生产级: 经 20+ 年大规模验证",
], icon="⚡")

add_card(slide, 7.2, 1.3, 5.8, 5.5, "本项目的优势", [
    "异常检测: EMA 行为基线 + 10 个插件检测器（Snort3 无）",
    "智能评分: 置信度 + 时段感知 + 自学习抑制（Snort3 无）",
    "攻击链: 基于 Kill Chain 的阶段映射（Snort3 无）",
    "GUI: 完整 tkinter GUI 内置（Snort3 无自带 GUI）",
    "热加载: 目录监控 3s 自动生效（Snort3 需 kill -HUP）",
    "白名单: GUI 即时操作（Snort3 需修改规则重载）",
    "报告: HTML + CSV + TXT 一键导出（Snort3 需 ELK 等）",
    "门槛: Python 脚本即改即生效（Snort3 需编译）",
], icon="⭐")


# ==================== Slide 11: 文件统计 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.3, 5, 0.7, "项目规模与测试体系", font_size=36, color=ACCENT, bold=True)

add_card(slide, 0.8, 1.3, 5.8, 2.5, "代码规模", [
    "总计 38 个文件，约 9,800 行 Python",
    "common/        6 文件  1,200 行  公共基础设施",
    "module1/       4 文件    800 行  抓包引擎+协议解析",
    "module2/       3 文件    900 行  特征匹配+热加载",
    "module3/      11 文件  2,200 行  异常检测(含10插件)",
    "module4/       1 文件  1,750 行  GUI 主窗口",
    "tests/         7 文件  1,500 行  单元+集成+演示",
], icon="📦")

add_card(slide, 7.2, 1.3, 5.8, 2.5, "测试覆盖", [
    "test_module1.py   模块1 单元测试（Scapy 包构造验证）",
    "test_module2.py   模块2 单元测试（Aho-Corasick + 规则）",
    "test_module3.py   模块3 单元测试（10个检测器）",
    "quick_test.py     快速集成测试",
    "live_demo.py      实机演示（模拟攻击 → GUI）",
    "check_env.py      环境检查（scapy+Npcap 状态）",
    "test_intelligence  S6/S7/S8 智能特性六场景验证",
], icon="🧪")

add_card(slide, 0.8, 4.2, 12.1, 2.8, "部署方式", [
    "方式A: 双击一键启动.bat → 菜单选择 GUI/抓包/演示/测试",
    "方式B: 双击 dist/NetworkAttackDetector.exe → 独立exe，无需Python",
    "方式C: python main.py                        → GUI 普通模式",
    "方式D: python main.py --auto                 → 自动检测（需Npcap+管理员）",
    "方式E: python main.py --pcap file.pcap --auto → 离线PCAP分析",
    "方式F: python tests/live_demo.py             → 模拟演示",
    "方式G: python test_intelligence.py           → 智能特性测试",
], icon="🚀")


# ==================== Slide 12: 后续规划 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.3, 5, 0.7, "后续改进方向", font_size=36, color=ACCENT, bold=True)

add_card(slide, 0.8, 1.3, 3.9, 2.5, "短期 (低成本高回报)", [
    "告警关联聚合:",
    "  同一IP多类型告警 →",
    "  Kill Chain 阶段推断",
    "",
    "C2 Beacon 检测:",
    "  连接间隔周期性分析",
    "",
    "配置档案切换:",
    "  home/office/dmz 多套",
], icon="📋")

add_card(slide, 5.1, 1.3, 3.9, 2.5, "中期 (架构增强)", [
    "告警输出插件化:",
    "  IAlertSink → Syslog",
    "  /Email/Webhook/ES",
    "",
    "数据源插件化:",
    "  IDataSource →",
    "  PCAPNG/Kafka/Zeek",
    "",
    "自动化响应:",
    "  Windows防火墙封禁",
    "  TCP RST 断开连接",
], icon="🔧")

add_card(slide, 9.4, 1.3, 3.9, 2.5, "长期 (技术演进)", [
    "机器学习引擎:",
    "  模块5 → Isolation",
    "  Forest / Autoencoder",
    "  无监督异常检测",
    "",
    "Web Dashboard:",
    "  Flask + ECharts",
    "  浏览器远程监控",
    "",
    "分布式部署:",
    "  Sensor + 汇聚中心",
    "  Kafka 消息队列",
], icon="🔮")

add_card(slide, 0.8, 4.2, 12.1, 2.8, "建议新增的 5 个检测器", [
    "1. C2 Beacon Detector      周期性心跳检测 → 识别 Cobalt Strike/Metasploit 等 C2 框架",
    "2. Payload Entropy Detector 载荷香农熵 >7.5 → 检测非标准端口加密恶意流量",
    "3. DNS Tunnel Detector      超长/高频 DNS 查询 → 检测 DNS 隧道数据外泄",
    "4. Alert Correlator         60s 窗口告警关联 → 碎片告警聚合为 Kill Chain 事件",
    "5. Protocol Masquerade      端口与协议不一致检测 → HTTP 跑非80端口 / TLS伪装",
], icon="🆕")


# ==================== Slide 13: 致谢 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 2, 2.0, 9, 1.2,
             "感谢聆听",
             font_size=52, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 2, 3.3, 9, 0.8,
             "网络攻击检测系统 (NIDS) v2.0",
             font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(
    1, Inches(4.5), Inches(4.3), Inches(4.3), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text_box(slide, 1.5, 4.7, 10, 1.5,
             "双引擎检测  |  10 个插件  |  智能评分  |  规则热加载  |  纯 Python 实现\n"
             "开源项目: github.com/gyhbwjhnb/web_attack_detection_system",
             font_size=14, color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER)


# ==================== 保存 ====================

os.makedirs("docs", exist_ok=True)
prs.save(OUTPUT)
print(f"PPT 已生成: {OUTPUT}")
print(f"共 {len(prs.slides)} 页幻灯片")
